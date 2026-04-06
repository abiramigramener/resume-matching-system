import re
from pathlib import Path
from typing import Dict, List
import numpy as np
from PIL import Image
import pytesseract
import easyocr
from docx import Document
import openpyxl
from pptx import Presentation
import fitz
from io import BytesIO
from fastapi import UploadFile
from app.shared.skills import SKILL_ALIASES

reader = easyocr.Reader(['en'], gpu=False)


def _normalize(text: str) -> str:
    text = text.lower()
    text = re.sub(r'[\n\r\t]+', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()


def _smart_title(s: str) -> str:
    """
    Title-case a skill string while preserving known acronyms.
    e.g. "pci" → "PCI", "ecg interpretation" → "ECG Interpretation"
    """
    ACRONYMS = {
        'pci', 'iabp', 'ecg', 'ekg', 'mri', 'ct', 'icu', 'ot', 'er',
        'sql', 'nosql', 'html', 'css', 'api', 'rest', 'aws', 'gcp',
        'ci/cd', 'nlp', 'llm', 'eda', 'etl', 'erp', 'sap', 'crm',
        'cad', 'bim', 'rcc', 'boq', 'plc', 'dcs', 'pcb', 'cnc',
        'seo', 'sem', 'smm', 'kpi', 'roi', 'b2b', 'b2c',
        'rn', 'mbbs', 'md', 'ms', 'phd',
    }
    words = s.split()
    result = []
    for w in words:
        if w.lower() in ACRONYMS:
            result.append(w.upper())
        else:
            result.append(w.capitalize())
    return ' '.join(result)


# Section headers that indicate a skills list follows
_SKILL_SECTION_HEADERS = re.compile(
    r'(?:^|\n)\s*(?:technical\s+)?(?:core\s+)?(?:key\s+)?skills?'
    r'(?:\s+(?:summary|set|profile|highlights?))?\s*[:\-–—]?\s*\n',
    re.IGNORECASE
)

# Sections that signal the skills block has ended
_NEXT_SECTION = re.compile(
    r'\n\s*(?:experience|education|employment|work history|projects?|'
    r'certifications?|awards?|publications?|references?|languages?|'
    r'interests?|hobbies|summary|objective|profile)\s*[:\-–—]?\s*\n',
    re.IGNORECASE
)


def _extract_from_skills_section(text: str) -> list[str]:
    """
    Parse items listed under a SKILLS / CORE SKILLS / KEY SKILLS section.
    Handles comma-separated, bullet-point, and multi-line formats.
    Returns smart-cased, deduplicated skill strings.
    """
    match = _SKILL_SECTION_HEADERS.search(text)
    if not match:
        return []

    block_start = match.end()
    next_sec = _NEXT_SECTION.search(text, block_start)
    block = text[block_start: next_sec.start() if next_sec else len(text)]

    # Split on commas, bullets, newlines, semicolons, pipes
    raw = re.split(r'[,;\n|]|[•·\*]|\d+\.', block)

    skills = []
    for item in raw:
        item = re.sub(r'\s+', ' ', item).strip().strip('–—:•·()-')
        if not item:
            continue
        if re.search(r'\d{4}', item):       # contains a year → likely a date/period
            continue
        if re.search(r'https?://', item):   # URL
            continue
        if len(item) > 50:                  # too long to be a skill name
            continue
        if item.isdigit():
            continue
        word_count = len(item.split())
        if word_count < 1 or word_count > 5:
            continue
        skills.append(_smart_title(item))

    return list(dict.fromkeys(skills))


def extract_skills(text: str) -> list[str]:
    """
    Hybrid skill extractor — returns deduplicated skill name strings.
    See extract_skill_details() for structured {name, years} output.
    """
    details = extract_skill_details(text)
    return [d["name"] for d in details]


# Patterns like "React - 2 years", "Node.js (3 yrs)", "5 years in Java"
_SKILL_EXP_PATTERNS = [
    re.compile(r'([A-Za-z][A-Za-z0-9\s\./\+\#\-]{1,40}?)\s*[\-–—:]\s*(\d+\.?\d*)\s*(?:years?|yrs?)', re.I),
    re.compile(r'([A-Za-z][A-Za-z0-9\s\./\+\#\-]{1,40}?)\s*\((\d+\.?\d*)\s*(?:years?|yrs?)\)', re.I),
    re.compile(r'(\d+\.?\d*)\s*(?:years?|yrs?)\s+(?:of\s+)?(?:experience\s+(?:in|with)\s+)?([A-Za-z][A-Za-z0-9\s\./\+\#\-]{1,40})', re.I),
]


def extract_skill_details(text: str) -> list[dict]:
    """
    Returns list of {name: str, years: float} dicts.

    Strategy:
    1. Try to extract per-skill experience from patterns in the text
    2. Fall back to section-based + alias extraction with years=0 (unknown)
    """
    # Try per-skill experience patterns first
    skill_years: dict[str, float] = {}
    for pat in _SKILL_EXP_PATTERNS:
        for m in pat.finditer(text):
            groups = m.groups()
            if groups[0].replace('.', '').replace(' ', '').isdigit():
                # Pattern 3: years first, skill second
                years_str, skill_str = groups[0], groups[1]
            else:
                skill_str, years_str = groups[0], groups[1]
            skill_name = _smart_title(skill_str.strip().strip('-–—:()'))
            if 2 <= len(skill_name) <= 50 and not skill_name[0].isdigit():
                try:
                    skill_years[skill_name] = float(years_str)
                except ValueError:
                    pass

    # Get full skill list via hybrid extraction
    section_skills = _extract_from_skills_section(text)
    normalized = _normalize(text)
    alias_skills = []
    for canonical, aliases in SKILL_ALIASES.items():
        for alias in aliases:
            pattern = r'(?<![a-zA-Z])' + re.escape(alias) + r'(?![a-zA-Z])'
            if re.search(pattern, normalized):
                alias_skills.append(_smart_title(canonical))
                break

    section_lower = {s.lower() for s in section_skills}
    all_skills = section_skills + [s for s in alias_skills if s.lower() not in section_lower]
    all_skills = list(dict.fromkeys(all_skills))

    # Merge: attach years where found, 0 where unknown
    result = []
    seen = set()
    for skill in all_skills:
        key = skill.lower()
        if key in seen:
            continue
        seen.add(key)
        # Try exact match, then partial match in skill_years
        years = skill_years.get(skill, 0.0)
        if years == 0.0:
            for k, v in skill_years.items():
                if k.lower() in key or key in k.lower():
                    years = v
                    break
        result.append({"name": skill, "years": years})

    return result


async def extract_text_from_file(file: UploadFile) -> str:
    if not file or not file.filename:
        return ""
    filename = file.filename.lower()
    ext = Path(filename).suffix
    text = ""
    try:
        await file.seek(0)
        contents = await file.read()
        if ext == '.pdf':
            text = _from_pdf(contents)
        elif ext in ['.doc', '.docx']:
            text = _from_docx(contents)
        elif ext in ['.xls', '.xlsx']:
            text = _from_excel(contents)
        elif ext in ['.ppt', '.pptx']:
            text = _from_pptx(contents)
        elif ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
            text = _from_image(contents)
        else:
            text = contents.decode('utf-8', errors='ignore')
    except Exception as e:
        print(f"Extraction error for {filename}: {e}")
        text = ""
    if ext == '.pdf' and len(text.strip()) < 50:
        try:
            await file.seek(0)
            contents = await file.read()
            text = _fallback_ocr(contents)
        except Exception as e:
            print(f"OCR fallback failed: {e}")
    return text.strip()


def parse_resume(text: str) -> Dict:
    """
    Extract structured fields from raw resume text.

    Returns
    -------
    skills              : list of canonical skill names
    experience_years    : total years of experience (float)
    education           : highest degree string (backward compat)
    education_details   : list of {degree, field, institution, year}
    job_title           : primary/most recent job title
    phone, email, name
    """
    skills = extract_skills(text)

    # ── Total Experience ──────────────────────────────────────────────────────
    experience = _extract_total_experience(text)

    # ── Contact ───────────────────────────────────────────────────────────────
    phone_match = re.search(r'(\+?\d[\d\s\-().]{7,}\d)', text)
    phone = phone_match.group(1).strip() if phone_match else None

    email_match = re.search(r'[\w.\-+]+@[\w\-]+\.[a-zA-Z]{2,}', text)
    email = email_match.group(0) if email_match else None

    # ── Name ──────────────────────────────────────────────────────────────────
    name = None
    skip = {"resume", "curriculum", "vitae", "cv", "profile", "summary",
            "objective", "contact", "details", "information"}
    for line in text.splitlines():
        line = line.strip()
        words = line.split()
        if (2 <= len(words) <= 4
                and all(re.match(r"^[a-zA-Z.\-']+$", w) for w in words)
                and not any(w.lower() in skip for w in words)
                and line[0].isupper()):
            name = line
            break

    # ── Education ─────────────────────────────────────────────────────────────
    edu_details = _extract_education(text)
    # Backward-compat flat string: highest degree
    education = edu_details[0]["degree"] if edu_details else "Not Specified"

    # ── Job Title ─────────────────────────────────────────────────────────────
    job_title = None
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    title_patterns = [
        r'^((?:senior|junior|lead|principal|staff|associate|chief|head of)?\s*'
        r'(?:software|data|machine learning|civil|mechanical|electrical|frontend|backend|'
        r'full.?stack|devops|cloud|ux|ui|product|project|marketing|sales|hr|finance|'
        r'medical|clinical|legal|research|business)\s*'
        r'(?:engineer|scientist|analyst|developer|designer|manager|consultant|'
        r'architect|specialist|officer|director|lead|intern))',
    ]
    for line in lines[:8]:
        for pat in title_patterns:
            m = re.match(pat, line, re.I)
            if m:
                job_title = m.group(1).strip().title()
                break
        if job_title:
            break

    return {
        "skills": skills,
        "experience_years": experience,
        "education": education,
        "education_details": edu_details,
        "job_title": job_title,
        "phone": phone,
        "email": email,
        "name": name,
    }


# ── Experience extraction helpers ─────────────────────────────────────────────

def _extract_total_experience(text: str) -> float:
    """
    Robustly extract total years of experience using multiple strategies:
    1. Explicit mention: "5 years of experience", "8+ years in software"
    2. Date range calculation: "2015 – 2023", "Jan 2018 to Present"
    3. Multiple role summation (best-effort)
    Returns 0.0 if nothing found.
    """
    # Strategy 1: explicit mention (most reliable)
    explicit_patterns = [
        r'(\d+\.?\d*)\+?\s*(?:years?|yrs?)\s+(?:of\s+)?(?:total\s+)?(?:work\s+)?experience',
        r'(?:over|more than|around|approximately|nearly)\s+(\d+\.?\d*)\s*(?:years?|yrs?)',
        r'(\d+\.?\d*)\+?\s*(?:years?|yrs?)\s+(?:in|of)\s+\w',
        r'experience\s+of\s+(\d+\.?\d*)\s*(?:years?|yrs?)',
        r'(\d+\.?\d*)\s*(?:years?|yrs?)\s+(?:of\s+)?(?:professional|industry|relevant|work)',
    ]
    for pat in explicit_patterns:
        m = re.search(pat, text, re.I)
        if m:
            try:
                val = float(m.group(1))
                if 0 < val <= 50:   # sanity check
                    return val
            except (ValueError, IndexError):
                pass

    # Strategy 2: calculate from date ranges in experience section
    years_from_dates = _calculate_experience_from_dates(text)
    if years_from_dates > 0:
        return years_from_dates

    # Strategy 3: fallback — any "N years" mention
    m = re.search(r'(\d+\.?\d*)\s*(?:years?|yrs?)', text, re.I)
    if m:
        try:
            val = float(m.group(1))
            if 0 < val <= 50:
                return val
        except ValueError:
            pass

    return 0.0


def _calculate_experience_from_dates(text: str) -> float:
    """
    Find date ranges in the experience section and sum up durations.
    Handles: "2015 – 2023", "Jan 2018 – Present", "2018 to 2021"
    """
    import datetime
    current_year = datetime.datetime.now().year

    # Find the experience/work section
    exp_section_match = re.search(
        r'\n\s*(?:experience|employment|work history|professional experience)\s*[:\-–—]?\s*\n',
        text, re.I
    )
    search_text = text[exp_section_match.start():] if exp_section_match else text

    # Pattern: year – year or year – present
    range_pattern = re.compile(
        r'(?:(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*\.?\s+)?'
        r'(\d{4})\s*(?:–|—|-|to)\s*'
        r'(?:(?:jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)[a-z]*\.?\s+)?'
        r'(\d{4}|present|current|now|till date|to date)',
        re.I
    )

    total_months = 0
    seen_ranges = set()
    for m in range_pattern.finditer(search_text):
        start_year = int(m.group(1))
        end_raw = m.group(2).lower()
        end_year = current_year if end_raw in ('present', 'current', 'now', 'till date', 'to date') else int(end_raw)

        if start_year < 1970 or start_year > current_year:
            continue
        if end_year < start_year or end_year > current_year + 1:
            continue

        key = (start_year, end_year)
        if key in seen_ranges:
            continue
        seen_ranges.add(key)
        total_months += (end_year - start_year) * 12

    if total_months > 0:
        # Cap at reasonable max and round to 1 decimal
        return round(min(total_months / 12, 50), 1)
    return 0.0


# ── Education extraction helpers ──────────────────────────────────────────────

# Degree aliases → canonical short form
_DEGREE_MAP = {
    r'\bb\.?\s*tech\.?\b': 'B.Tech',
    r'\bm\.?\s*tech\.?\b': 'M.Tech',
    r'\bb\.?\s*e\.?\b': 'B.E',
    r'\bm\.?\s*e\.?\b': 'M.E',
    r'\bb\.?\s*sc\.?\b': 'B.Sc',
    r'\bm\.?\s*sc\.?\b': 'M.Sc',
    r'\bb\.?\s*com\.?\b': 'B.Com',
    r'\bm\.?\s*com\.?\b': 'M.Com',
    r'\bb\.?\s*b\.?\s*a\.?\b': 'BBA',
    r'\bm\.?\s*b\.?\s*a\.?\b': 'MBA',
    r'\bb\.?\s*c\.?\s*a\.?\b': 'BCA',
    r'\bm\.?\s*c\.?\s*a\.?\b': 'MCA',
    r'\bmbbs\b': 'MBBS',
    r'\bm\.?\s*d\.?\b': 'MD',
    r'\bd\.?\s*m\.?\b': 'DM',
    r'\bm\.?\s*s\.?\b': 'MS',
    r'\bph\.?\s*d\.?\b': 'PhD',
    r'\bbachelor(?:\'s)?\s+of\s+(?:science|engineering|technology|arts|commerce|business)': 'Bachelor',
    r'\bmaster(?:\'s)?\s+of\s+(?:science|engineering|technology|arts|commerce|business|administration)': 'Master',
    r'\bdoctor\s+of\s+philosophy\b': 'PhD',
    r'\bdiploma\b': 'Diploma',
    r'\bhigh\s+school\b|\bssc\b|\bhsc\b|\b10th\b|\b12th\b': 'High School',
}

# Common fields of study — matches "in Computer Science", "Computer Science", "- Civil Engineering"
_FIELD_PATTERN = re.compile(
    r'(?:(?:in|of)\s+)?([A-Z][a-zA-Z\s&/]{3,40}?)(?:\s*[,|\(|\n]|\s*\d{4}|$)',
    re.I
)

# Year pattern
_YEAR_PATTERN = re.compile(r'\b(19[89]\d|20[012]\d)\b')


def _extract_education(text: str) -> list[dict]:
    """
    Extract structured education entries from resume text.
    Returns list of {degree, field, institution, year}.
    """
    # Find education section
    edu_match = re.search(
        r'\n\s*(?:education(?:al)?\s+(?:qualification|background|details)?|'
        r'academic\s+(?:qualification|background|details)?|qualifications?)\s*[:\-–—]?\s*\n',
        text, re.I
    )
    # Search in education section if found, else full text
    search_text = text[edu_match.start():] if edu_match else text

    # Find next major section after education
    next_sec = re.search(
        r'\n\s*(?:experience|employment|skills?|projects?|certifications?|work)\s*[:\-–—]?\s*\n',
        search_text[50:] if len(search_text) > 50 else search_text, re.I
    )
    if next_sec:
        search_text = search_text[:next_sec.start() + 50]

    entries = []
    seen_degrees = set()

    for deg_pattern, canonical in _DEGREE_MAP.items():
        for m in re.finditer(deg_pattern, search_text, re.I):
            if canonical in seen_degrees:
                continue
            seen_degrees.add(canonical)

            # Get surrounding context (100 chars after match)
            ctx_start = m.start()
            ctx_end = min(m.end() + 150, len(search_text))
            ctx = search_text[ctx_start:ctx_end]

            # Extract field of study
            # Try "in/of <field>" first, then field directly after degree on same line
            field = None
            # Get the line containing the degree match
            line_start = search_text.rfind('\n', 0, m.start()) + 1
            line_end = search_text.find('\n', m.end())
            if line_end == -1:
                line_end = len(search_text)
            degree_line = search_text[line_start:line_end]

            # Pattern 1: "B.Tech in Computer Science" or "B.E. Computer Science"
            after_degree = degree_line[m.end() - line_start:].strip().lstrip('.,- ')
            # Remove year and institution keywords from the tail
            after_degree = re.sub(r'\b\d{4}\b.*', '', after_degree)
            after_degree = re.sub(r'(?:university|college|institute|iit|nit|from|aiims).*', '', after_degree, flags=re.I)
            after_degree = re.sub(r'^(?:in|of)\s+', '', after_degree, flags=re.I).strip().strip('.,|-')

            if after_degree and 2 <= len(after_degree.split()) <= 5 and len(after_degree) >= 3:
                field = after_degree.title()

            # Extract year
            year = None
            ym = _YEAR_PATTERN.search(ctx)
            if ym:
                year = int(ym.group(1))

            # Extract institution — look for lines containing institution keywords
            institution = None
            inst_pattern = re.compile(
                r'\b([A-Z][A-Za-z\s,\.]{3,50}?'
                r'(?:university|college|institute|school|iit|nit|bits|iiit|iim|aiims))\b',
                re.I
            )
            # Search only after the degree match to avoid capturing the degree itself
            post_degree = search_text[m.end():m.end() + 200]
            im = inst_pattern.search(post_degree)
            if im:
                institution = im.group(1).strip().title()

            entries.append({
                "degree": canonical,
                "field": field,
                "institution": institution,
                "year": year,
            })

    # Sort by degree seniority (PhD > Master > Bachelor > etc.)
    _DEGREE_ORDER = ['PhD', 'DM', 'MD', 'MS', 'Master', 'M.Tech', 'M.E', 'M.Sc', 'MCA', 'MBA', 'M.Com',
                     'MBBS', 'Bachelor', 'B.Tech', 'B.E', 'B.Sc', 'BCA', 'BBA', 'B.Com', 'Diploma', 'High School']
    entries.sort(key=lambda e: _DEGREE_ORDER.index(e['degree']) if e['degree'] in _DEGREE_ORDER else 99)

    return entries if entries else []


def _from_pdf(contents: bytes) -> str:
    doc = fitz.open(stream=contents, filetype="pdf")
    text = ""
    for page in doc:
        page_text = page.get_text("text")
        if page_text.strip():
            text += page_text
        else:
            pix = page.get_pixmap()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text += pytesseract.image_to_string(img) + "\n"
    doc.close()
    return text


def _from_docx(contents: bytes) -> str:
    return "\n".join([p.text for p in Document(BytesIO(contents)).paragraphs])


def _from_excel(contents: bytes) -> str:
    wb = openpyxl.load_workbook(BytesIO(contents), read_only=True)
    text = ""
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            text += " | ".join([str(c) if c is not None else "" for c in row]) + "\n"
    return text


def _from_pptx(contents: bytes) -> str:
    prs = Presentation(BytesIO(contents))
    return "\n".join(shape.text for slide in prs.slides
                     for shape in slide.shapes if hasattr(shape, "text"))


def _from_image(contents: bytes) -> str:
    img = Image.open(BytesIO(contents))
    text = pytesseract.image_to_string(img)
    if len(text.strip()) < 30:
        text = " ".join([r[1] for r in reader.readtext(np.array(img))])
    return text


def _fallback_ocr(contents: bytes) -> str:
    try:
        return pytesseract.image_to_string(Image.open(BytesIO(contents)))
    except Exception:
        return ""
