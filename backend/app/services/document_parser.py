import re
from pathlib import Path
from typing import Dict, List
import numpy as np
from PIL import Image
import pytesseract
import easyocr
from docx import Document
import openpyxl
from pptx import Presentation
import fitz
from io import BytesIO
from fastapi import UploadFile

reader = easyocr.Reader(['en'], gpu=False)

SKILL_ALIASES = {
    "machine learning":     ["machine learning", "ml model", "supervised learning", "unsupervised learning"],
    "deep learning":        ["deep learning", "neural network", "neural networks"],
    "nlp":                  ["nlp", "natural language processing", "text mining", "text analytics"],
    "computer vision":      ["computer vision", "image recognition", "object detection"],
    "data science":         ["data science", "data scientist"],
    "data analysis":        ["data analysis", "data analyst", "data analytics", "eda", "exploratory data analysis"],
    "data visualization":   ["data visualization", "data viz", "tableau", "power bi", "powerbi", "looker", "matplotlib", "seaborn", "plotly"],
    "statistics":           ["statistics", "statistical analysis", "statistical modeling", "hypothesis testing"],
    "feature engineering":  ["feature engineering", "feature selection"],
    "model deployment":     ["model deployment", "mlops", "model serving"],
    "time series":          ["time series", "forecasting", "arima"],
    "llm":                  ["llm", "large language model", "gpt", "bert", "transformers", "huggingface"],
    "a/b testing":          ["a/b testing", "ab testing", "split testing"],
    "python":               ["python"],
    "pandas":               ["pandas"],
    "numpy":                ["numpy"],
    "scikit-learn":         ["scikit-learn", "sklearn", "scikit learn"],
    "tensorflow":           ["tensorflow"],
    "keras":                ["keras"],
    "pytorch":              ["pytorch"],
    "xgboost":              ["xgboost", "lightgbm", "gradient boosting"],
    "scipy":                ["scipy"],
    "jupyter":              ["jupyter"],
    "sql":                  ["sql", "mysql", "postgresql", "postgres", "sqlite"],
    "nosql":                ["mongodb", "cassandra", "dynamodb", "nosql"],
    "spark":                ["spark", "pyspark", "apache spark"],
    "hadoop":               ["hadoop", "hive", "hdfs"],
    "kafka":                ["kafka"],
    "airflow":              ["airflow", "luigi"],
    "aws":                  ["aws", "amazon web services", "sagemaker"],
    "azure":                ["azure", "microsoft azure"],
    "gcp":                  ["gcp", "google cloud", "bigquery", "dataflow"],
    "docker":               ["docker"],
    "kubernetes":           ["kubernetes", "k8s"],
    "git":                  ["git", "github", "gitlab"],
    "linux":                ["linux", "unix", "bash", "shell scripting"],
    "ci/cd":                ["ci/cd", "jenkins", "github actions"],
    "javascript":           ["javascript"],
    "typescript":           ["typescript"],
    "react":                ["react", "reactjs"],
    "vue":                  ["vue", "vuejs"],
    "angular":              ["angular"],
    "java":                 ["java"],
    "scala":                ["scala"],
    "r":                    ["r programming", "rstudio", "tidyverse"],
    "excel":                ["excel", "spreadsheet"],
    "power bi":             ["power bi", "powerbi"],
    "tableau":              ["tableau"],
    "api":                  ["rest api", "restful api", "fastapi", "flask", "django"],
    "devops":               ["devops"],
    "mlflow":               ["mlflow"],
    "reinforcement learning": ["reinforcement learning"],
}


def _normalize(text):
    text = text.lower()
    text = re.sub(r'[\n\r\t]+', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()


def extract_skills(text):
    normalized = _normalize(text)
    found = []
    for canonical, aliases in SKILL_ALIASES.items():
        for alias in aliases:
            pattern = r'(?<![a-zA-Z])' + re.escape(alias) + r'(?![a-zA-Z])'
            if re.search(pattern, normalized):
                found.append(canonical)
                break
    return list(dict.fromkeys(found))


async def extract_text_from_file(file):
    if not file or not file.filename:
        return ""
    filename = file.filename.lower()
    ext = Path(filename).suffix
    text = ""
    try:
        await file.seek(0)
        contents = await file.read()
        if ext == '.pdf':
            text = extract_from_pdf(contents)
        elif ext in ['.doc', '.docx']:
            text = extract_from_docx(contents)
        elif ext in ['.xls', '.xlsx']:
            text = extract_from_excel(contents)
        elif ext in ['.ppt', '.pptx']:
            text = extract_from_pptx(contents)
        elif ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
            text = extract_from_image(contents)
        else:
            text = contents.decode('utf-8', errors='ignore')
    except Exception as e:
        print(f"Extraction error for {filename}: {e}")
        text = ""
    if ext == '.pdf' and len(text.strip()) < 50:
        try:
            await file.seek(0)
            contents = await file.read()
            text = fallback_ocr(contents)
        except Exception as ocr_err:
            print(f"OCR fallback failed: {ocr_err}")
    return text.strip()


def extract_from_pdf(contents):
    doc = fitz.open(stream=contents, filetype="pdf")
    text = ""
    for page in doc:
        page_text = page.get_text("text")
        if page_text.strip():
            text += page_text
        else:
            pix = page.get_pixmap()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text += pytesseract.image_to_string(img) + "\n"
    doc.close()
    return text


def extract_from_docx(contents):
    doc = Document(BytesIO(contents))
    return "\n".join([para.text for para in doc.paragraphs])


def extract_from_excel(contents):
    wb = openpyxl.load_workbook(BytesIO(contents), read_only=True)
    text = ""
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            text += " | ".join([str(c) if c is not None else "" for c in row]) + "\n"
    return text


def extract_from_pptx(contents):
    prs = Presentation(BytesIO(contents))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_from_image(contents):
    img = Image.open(BytesIO(contents))
    text = pytesseract.image_to_string(img)
    if len(text.strip()) < 30:
        result = reader.readtext(np.array(img))
        text = " ".join([res[1] for res in result])
    return text


def fallback_ocr(contents):
    try:
        img = Image.open(BytesIO(contents))
        return pytesseract.image_to_string(img)
    except Exception:
        return ""


def parse_structured_data(text):
    skills = extract_skills(text)
    exp_patterns = [
        r'(\d+)\+?\s*(?:years?|yrs?)\s*(?:of)?\s*(?:experience|exp)',
        r'(\d+)\s*(?:years?|yrs?)\s*(?:exp|experience)',
    ]
    experience = 0
    for pattern in exp_patterns:
        m = re.search(pattern, text, re.I)
        if m:
            experience = int(m.group(1))
            break
    phone_match = re.search(r'(\+?\d[\d\s\-().]{7,}\d)', text)
    phone = phone_match.group(1).strip() if phone_match else None
    email_match = re.search(r'[\w.\-+]+@[\w\-]+\.[a-zA-Z]{2,}', text)
    email = email_match.group(0) if email_match else None
    name = None
    skip_words = {"resume", "curriculum", "vitae", "cv", "profile", "summary",
                  "objective", "contact", "details", "information"}
    for line in text.splitlines():
        line = line.strip()
        words = line.split()
        if (2 <= len(words) <= 4
                and all(re.match(r"^[a-zA-Z.\-']+$", w) for w in words)
                and not any(w.lower() in skip_words for w in words)
                and line[0].isupper()):
            name = line
            break
    edu_keywords = ["bachelor", "master", "phd", "b.tech", "m.tech", "bsc", "msc", "degree"]
    education = next((kw.capitalize() for kw in edu_keywords if kw in text.lower()), "Not Specified")
    return {
        "skills": skills,
        "experience_years": experience,
        "education": education,
        "phone": phone,
        "email": email,
        "name": name,
    }
